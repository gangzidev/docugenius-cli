#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Document converter for DocuGenius CLI
Supports Word, Excel, PowerPoint, and PDF conversion to Markdown
"""

import sys
import logging
import re
from pathlib import Path
from typing import Optional, List
from dataclasses import dataclass

from .config import Config
from .image_extractor import ImageExtractor

logger = logging.getLogger(__name__)

# Markdown escaping regex patterns
_RE_COLLAPSE_WHITESPACE = re.compile(r"\s+")
_RE_COLLAPSE_EXTRA_BLANK_LINES = re.compile(r"\n{3,}")
_RE_ESCAPE_MARKDOWN_LEADING = re.compile(r"^([>#\-\+\*])")
_RE_ESCAPE_MARKDOWN_ORDERED_LIST = re.compile(r"^(\d+)\.\s")


@dataclass
class ConversionResult:
    """Result of a document conversion."""

    success: bool
    output_path: Optional[Path] = None
    images_extracted: int = 0
    images_dir: Optional[Path] = None
    error: Optional[str] = None


class Converter:
    """Document converter supporting multiple formats."""

    def __init__(self, config: Config):
        """
        Initialize converter.

        Args:
            config: Configuration object
        """
        self.config = config
        self.image_extractor = ImageExtractor(config)

    def convert_file(
        self,
        file_path: Path,
        extract_images: bool = False,
        split_threshold: int = 500000,
        overwrite: bool = False,
        verbose: bool = False,
    ) -> ConversionResult:
        """
        Convert a document file to Markdown.

        Args:
            file_path: Path to the input file
            extract_images: Whether to extract images
            split_threshold: Character count threshold for splitting documents
            overwrite: Whether to overwrite existing output files
            verbose: Enable verbose logging

        Returns:
            ConversionResult object with conversion status and output path
        """
        try:
            # Validate file exists
            if not file_path.exists():
                return ConversionResult(success=False, error=f"File not found: {file_path}")

            # Get file extension
            ext = file_path.suffix.lower()

            # Determine output path
            output_dir = Path(self.config.output_dir)
            if not output_dir.is_absolute():
                output_dir = file_path.parent / output_dir
            output_dir.mkdir(parents=True, exist_ok=True)

            output_path = output_dir / f"{file_path.stem}.md"

            # Check if output exists and overwrite is disabled
            if output_path.exists() and not overwrite:
                if verbose:
                    logger.info(f"Output file already exists, skipping: {output_path}")
                return ConversionResult(success=True, output_path=output_path)

            # Convert based on file type
            if ext == ".docx":
                markdown_content = self._convert_docx(file_path)
            elif ext == ".xlsx":
                markdown_content = self._convert_xlsx(file_path)
            elif ext == ".pptx":
                markdown_content = self._convert_pptx(file_path)
            elif ext == ".pdf":
                markdown_content = self._convert_pdf(file_path)
            else:
                return ConversionResult(success=False, error=f"Unsupported file type: {ext}")

            # Check if content is empty
            if not markdown_content or len(markdown_content.strip()) == 0:
                return ConversionResult(
                    success=False,
                    error="No content could be extracted from the document. The file may be image-only or use an unsupported format."
                )

            # Extract images if requested
            images_extracted = 0
            images_dir = None
            if extract_images:
                result = self.image_extractor.extract_images(file_path, output_dir)
                images_extracted = result.count
                images_dir = result.output_dir
                if result.markdown_content:
                    markdown_content = result.markdown_content

            # Split document if too large
            if len(markdown_content) > split_threshold:
                self._split_document(markdown_content, output_path, file_path.name)
            else:
                # Write output file
                output_path.write_text(markdown_content, encoding="utf-8")

            return ConversionResult(
                success=True,
                output_path=output_path,
                images_extracted=images_extracted,
                images_dir=images_dir,
            )

        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}", exc_info=True)
            return ConversionResult(success=False, error=str(e))

    def _convert_docx(self, file_path: Path) -> str:
        """Convert Word document to Markdown."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for Word conversion. Install with: pip install python-docx")

        doc = Document(file_path)
        markdown_lines = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Handle headings
            if para.style.name.startswith("Heading"):
                level = int(para.style.name.replace("Heading ", ""))
                markdown_lines.append(f"{'#' * level} {text}")
            else:
                # Escape markdown special characters
                text = self._escape_markdown(text)
                markdown_lines.append(text)

        # Handle tables
        for table in doc.tables:
            markdown_lines.append("")
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                markdown_lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    markdown_lines.append("|" + "|".join(["---" for _ in cells]) + "|")
            markdown_lines.append("")

        return "\n".join(markdown_lines)

    def _convert_xlsx(self, file_path: Path) -> str:
        """Convert Excel spreadsheet to Markdown."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required for Excel conversion. Install with: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path, data_only=True)
        markdown_lines = [f"# {file_path.name}\n"]

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            markdown_lines.append(f"\n## {sheet_name}\n")

            # Find the actual data range
            data = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    data.append(row)

            if not data:
                continue

            # Convert to Markdown table
            for i, row in enumerate(data):
                cells = [str(cell) if cell is not None else "" for cell in row]
                markdown_lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    markdown_lines.append("|" + "|".join(["---" for _ in cells]) + "|")

            markdown_lines.append("")

        return "\n".join(markdown_lines)

    def _convert_pptx(self, file_path: Path) -> str:
        """Convert PowerPoint presentation to Markdown."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint conversion. Install with: pip install python-pptx"
            )

        prs = Presentation(file_path)
        markdown_lines = [f"# {file_path.name}\n"]

        for slide_idx, slide in enumerate(prs.slides, 1):
            markdown_lines.append(f"\n## Slide {slide_idx}\n")

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    text = self._escape_markdown(text)
                    markdown_lines.append(text)

            # Extract tables
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for i, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        markdown_lines.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            markdown_lines.append("|" + "|".join(["---" for _ in cells]) + "|")
                    markdown_lines.append("")

        return "\n".join(markdown_lines)

    def _convert_pdf(self, file_path: Path) -> str:
        """Convert PDF document to Markdown."""
        try:
            import pdfplumber
        except ImportError:
            raise ImportError(
                "pdfplumber is required for PDF conversion. Install with: pip install pdfplumber"
            )

        markdown_lines = [f"# {file_path.name}\n"]

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    markdown_lines.append(f"\n## Page {page_num}\n")
                    markdown_lines.append(text)

        return "\n".join(markdown_lines)

    def _escape_markdown(self, text: str) -> str:
        """Escape markdown special characters in text."""
        # Escape leading characters that would be interpreted as markdown
        lines = text.split("\n")
        escaped_lines = []

        for line in lines:
            # Escape ordered list markers (e.g., "1. " -> "1\. ")
            if _RE_ESCAPE_MARKDOWN_ORDERED_LIST.match(line):
                line = re.sub(r'^(\d+)\.', r'\1\.', line)

            # Escape other markdown characters at line start
            line = _RE_ESCAPE_MARKDOWN_LEADING.sub(r"\\\1", line)

            escaped_lines.append(line)

        return "\n".join(escaped_lines)

    def _split_document(self, content: str, output_path: Path, filename: str) -> None:
        """
        Split a large document into multiple files.

        Args:
            content: Full markdown content
            output_path: Base output path (without _part suffix)
            filename: Original filename for headers
        """
        # Split by headers (## or ###)
        parts = []
        current_part = []
        current_size = 0

        lines = content.split("\n")
        for line in lines:
            # Check if this is a section header
            if line.startswith("##") and current_size > self.config.split_threshold:
                parts.append("\n".join(current_part))
                current_part = [line]
                current_size = len(line)
            else:
                current_part.append(line)
                current_size += len(line) + 1

        # Add the last part
        if current_part:
            parts.append("\n".join(current_part))

        # If only one part, save as single file
        if len(parts) == 1:
            output_path.write_text(parts[0], encoding="utf-8")
            return

        # Save multiple parts
        base_path = output_path.parent / output_path.stem
        for i, part in enumerate(parts, 1):
            part_path = base_path.parent / f"{base_path.name}_part{i}.md"
            part_content = f"# {filename} - Part {i} of {len(parts)}\n\n{part}"
            part_path.write_text(part_content, encoding="utf-8")

        # Create index file
        index_content = f"# {filename} - Index\n\n"
        index_content += "This document has been split into multiple parts:\n\n"
        for i in range(1, len(parts) + 1):
            index_content += f"- [Part {i}]({base_path.name}_part{i}.md)\n"

        index_path = base_path.parent / f"{base_path.name}_index.md"
        index_path.write_text(index_content, encoding="utf-8")
