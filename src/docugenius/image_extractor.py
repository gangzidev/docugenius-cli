#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Image extractor for DocuGenius CLI
Extracts images from Word, PowerPoint, and PDF documents
"""

import logging
import re
import struct
from pathlib import Path
from typing import Optional, List
from dataclasses import dataclass

logger = logging.getLogger(__name__)

# Image format signatures (magic numbers)
_IMAGE_SIGNATURES = {
    b"\x89PNG\r\n\x1a\n": "png",
    b"\xff\xd8\xff": "jpeg",
    b"GIF87a": "gif",
    b"GIF89a": "gif",
    b"BM": "bmp",
    b"II\x2a\x00": "tiff",
    b"MM\x00\x2a": "tiff",
    b"\xd7\xcd\xc6\x9a": "wmf",
    b"\x01\x00\x00\x00": "emf",
}

# Image extraction constants
MIN_IMAGE_DIMENSION_PX = 20  # Images smaller than this are considered decorative
MAX_ASPECT_RATIO = 10.0  # Images with aspect ratio > this are considered decorative lines
MIN_IMAGE_DATA_BYTES = 500  # Images with less data than this are considered solid/transparent placeholders


@dataclass
class ImageExtractionResult:
    """Result of image extraction."""

    count: int
    output_dir: Optional[Path]
    markdown_content: Optional[str] = None


class ImageExtractor:
    """Extract images from documents."""

    def __init__(self, config):
        """
        Initialize image extractor.

        Args:
            config: Configuration object
        """
        self.config = config

    def extract_images(self, file_path: Path, output_dir: Path) -> ImageExtractionResult:
        """
        Extract images from a document.

        Args:
            file_path: Path to the input file
            output_dir: Directory to save extracted images

        Returns:
            ImageExtractionResult with count and output directory
        """
        ext = file_path.suffix.lower()

        try:
            if ext == ".docx":
                return self._extract_from_docx(file_path, output_dir)
            elif ext == ".pptx":
                return self._extract_from_pptx(file_path, output_dir)
            elif ext == ".pdf":
                return self._extract_from_pdf(file_path, output_dir)
            else:
                return ImageExtractionResult(count=0, output_dir=None)
        except Exception as e:
            logger.error(f"Error extracting images from {file_path}: {e}", exc_info=True)
            return ImageExtractionResult(count=0, output_dir=None)

    def _extract_from_docx(self, file_path: Path, output_dir: Path) -> ImageExtractionResult:
        """Extract images from Word document."""
        from docx import Document
        from docx.oxml import parse_xml
        from docx.oxml.ns import qn

        doc = Document(file_path)
        images_dir = output_dir / "images" / file_path.stem
        images_dir.mkdir(parents=True, exist_ok=True)

        image_count = 0
        image_references = []

        # Extract images from document
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    image_format = self._detect_image_format(image_data)

                    if image_format and self._is_valid_image(image_data):
                        image_name = f"image_{image_count + 1}.{image_format}"
                        image_path = images_dir / image_name

                        with open(image_path, "wb") as f:
                            f.write(image_data)

                        image_count += 1
                        image_references.append(image_name)
                except Exception as e:
                    logger.warning(f"Failed to extract image: {e}")
                    continue

        return ImageExtractionResult(count=image_count, output_dir=images_dir)

    def _extract_from_pptx(self, file_path: Path, output_dir: Path) -> ImageExtractionResult:
        """Extract images from PowerPoint presentation."""
        from pptx import Presentation

        prs = Presentation(file_path)
        images_dir = output_dir / "images" / file_path.stem
        images_dir.mkdir(parents=True, exist_ok=True)

        image_count = 0

        # Extract images from all slides
        for slide_idx, slide in enumerate(prs.slides, 1):
            for rel in slide.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        image_format = self._detect_image_format(image_data)

                        if image_format and self._is_valid_image(image_data):
                            image_name = f"slide{slide_idx}_image_{image_count + 1}.{image_format}"
                            image_path = images_dir / image_name

                            with open(image_path, "wb") as f:
                                f.write(image_data)

                            image_count += 1
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide {slide_idx}: {e}")
                        continue

        return ImageExtractionResult(count=image_count, output_dir=images_dir)

    def _extract_from_pdf(self, file_path: Path, output_dir: Path) -> ImageExtractionResult:
        """Extract images from PDF document."""
        try:
            import pdfplumber
        except ImportError:
            logger.warning("pdfplumber not available for image extraction")
            return ImageExtractionResult(count=0, output_dir=None)

        images_dir = output_dir / "images" / file_path.stem
        images_dir.mkdir(parents=True, exist_ok=True)

        image_count = 0

        try:
            with pdfplumber.open(file_path) as pdf:
                for page_idx, page in enumerate(pdf.pages, 1):
                    if page.images:
                        for img_idx, img in enumerate(page.images, 1):
                            try:
                                # Get image from PDF
                                image_data = page.within_bbox((img["x0"], img["top"], img["x1"], img["bottom"])).to_image()

                                # Convert to bytes
                                import io
                                from PIL import Image

                                img_bytes = io.BytesIO()
                                image_data.save(img_bytes, format="PNG")
                                img_bytes.seek(0)
                                image_data_bytes = img_bytes.read()

                                if self._is_valid_image(image_data_bytes):
                                    image_name = f"page{page_idx}_image{img_idx}.png"
                                    image_path = images_dir / image_name

                                    with open(image_path, "wb") as f:
                                        f.write(image_data_bytes)

                                    image_count += 1
                            except Exception as e:
                                logger.warning(f"Failed to extract image from page {page_idx}: {e}")
                                continue
        except Exception as e:
            logger.error(f"Error extracting images from PDF: {e}")

        return ImageExtractionResult(count=image_count, output_dir=images_dir)

    def _detect_image_format(self, image_data: bytes) -> Optional[str]:
        """
        Detect image format from binary data.

        Args:
            image_data: Binary image data

        Returns:
            Image format string (e.g., 'png', 'jpeg') or None if unknown
        """
        for signature, format_name in _IMAGE_SIGNATURES.items():
            if image_data.startswith(signature):
                return format_name
        return None

    def _is_valid_image(self, image_data: bytes) -> bool:
        """
        Check if image is valid (not decorative).

        Args:
            image_data: Binary image data

        Returns:
            True if image is valid, False otherwise
        """
        # Check minimum data size
        if len(image_data) < MIN_IMAGE_DATA_BYTES:
            return False

        # Try to load with PIL to get dimensions
        try:
            from PIL import Image
            import io

            img = Image.open(io.BytesIO(image_data))
            width, height = img.size

            # Check minimum dimensions
            if width < MIN_IMAGE_DIMENSION_PX or height < MIN_IMAGE_DIMENSION_PX:
                return False

            # Check aspect ratio
            aspect_ratio = max(width, height) / min(width, height)
            if aspect_ratio > MAX_ASPECT_RATIO:
                return False

            return True
        except Exception:
            # If we can't load the image, assume it's invalid
            return False
